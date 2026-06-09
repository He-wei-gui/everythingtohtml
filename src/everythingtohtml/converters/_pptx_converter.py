"""PPTX -> HTML, one ``<section>`` per slide (requires the ``pptx`` extra)."""

from __future__ import annotations

import base64
from typing import Any, BinaryIO

from .._base_converter import DocumentConverter, DocumentConverterResult
from .._exceptions import MissingDependencyException
from .._html_builder import escape_attr, escape_text, wrap_document
from .._stream_info import StreamInfo

__all__ = ["PptxConverter"]

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

_PPTX_STYLE = """
body { max-width: 72rem; }
.slide {
  position: relative;
  width: 100%;
  aspect-ratio: var(--pptx-slide-width) / var(--pptx-slide-height);
  margin: 1.5rem 0 2.5rem;
  overflow: hidden;
  border: 1px solid #d9d9df;
  border-radius: 8px;
  background: #fff;
  color: #111;
  box-shadow: 0 1px 6px rgb(0 0 0 / 0.08);
}
.slide-number {
  margin: 0 0 0.5rem;
  color: #666;
  font-size: 0.85rem;
}
.pptx-shape {
  position: absolute;
  overflow: hidden;
}
.pptx-text {
  line-height: 1.2;
  overflow-wrap: anywhere;
  white-space: normal;
}
.pptx-text p {
  margin: 0 0 0.25em;
}
.pptx-text p.bullet::before {
  content: "• ";
}
.pptx-image {
  display: block;
  width: 100%;
  height: 100%;
  object-fit: contain;
}
.pptx-table {
  width: 100%;
  height: 100%;
  margin: 0;
  table-layout: fixed;
  font-size: 0.8rem;
  background: #fff;
}
.pptx-table td,
.pptx-table th {
  padding: 0.2rem 0.3rem;
  vertical-align: top;
}
"""


class PptxConverter(DocumentConverter):
    """Convert PowerPoint ``.pptx`` decks into a slide-per-section HTML document."""

    priority = DocumentConverter.PRIORITY_SPECIFIC_FILE_FORMAT

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> bool:
        ext = stream_info.normalized_extension()
        mimetype = (stream_info.mimetype or "").split(";", 1)[0].strip().lower()
        return ext == ".pptx" or mimetype == _PPTX_MIME

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> DocumentConverterResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise MissingDependencyException(
                "PPTX support requires the 'pptx' extra. "
                "Install it with: pip install everythingtohtml[pptx]"
            ) from exc

        presentation = Presentation(file_stream)
        parts: list[str] = []
        title: str | None = None

        slide_width = _emu(presentation.slide_width)
        slide_height = _emu(presentation.slide_height)

        for index, slide in enumerate(presentation.slides, start=1):
            parts.append(f'<p class="slide-number">Slide {index}</p>')
            parts.append(
                f'<section class="slide" id="slide-{index}" '
                f'style="--pptx-slide-width:{slide_width};--pptx-slide-height:{slide_height};">'
            )
            for shape in slide.shapes:
                parts.append(_render_shape(shape, slide_width, slide_height))
                if title is None:
                    title = _shape_title(shape)
            parts.append("</section>")

        title = title or stream_info.filename
        html = wrap_document(
            "\n".join(parts),
            title=title,
            extra_head=f"<style>{_PPTX_STYLE}</style>\n",
        )
        return DocumentConverterResult(html, title=title)


def _render_shape(shape: Any, slide_width: int, slide_height: int) -> str:
    style = _shape_style(shape, slide_width, slide_height)
    body = _render_shape_body(shape)
    if not body:
        return ""
    return f'<div class="pptx-shape" style="{escape_attr(style)}">{body}</div>'


def _render_shape_body(shape: Any) -> str:
    image = getattr(shape, "image", None)
    if image is not None:
        return _render_image(image)
    if getattr(shape, "has_table", False):
        return _render_table(shape.table)
    if not getattr(shape, "has_text_frame", False):
        return ""
    blocks: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        html = "".join(_render_run(run) for run in paragraph.runs).strip()
        if not html:
            html = escape_text(paragraph.text.strip())
        if not html:
            continue
        level = paragraph.level or 0
        class_name = ' class="bullet"' if level > 0 else ""
        indent = f' style="padding-left:{level * 1.2:.1f}em"' if level > 0 else ""
        blocks.append(f"<p{class_name}{indent}>{html}</p>")
    return f'<div class="pptx-text">{"".join(blocks)}</div>' if blocks else ""


def _render_run(run: Any) -> str:
    text = run.text
    if not text:
        return ""
    content = escape_text(text).replace("\n", "<br>")
    style = _run_style(run)
    if style:
        content = f'<span style="{escape_attr(style)}">{content}</span>'
    font = run.font
    if font.bold:
        content = f"<strong>{content}</strong>"
    if font.italic:
        content = f"<em>{content}</em>"
    return content


def _run_style(run: Any) -> str:
    styles: list[str] = []
    font = run.font
    if font.size is not None:
        styles.append(f"font-size:{font.size.pt:.2f}pt")
    if font.underline:
        styles.append("text-decoration:underline")
    return ";".join(styles)


def _render_image(image: Any) -> str:
    content_type = image.content_type or f"image/{image.ext}"
    encoded = base64.b64encode(image.blob).decode("ascii")
    alt = image.filename or "slide image"
    return (
        f'<img class="pptx-image" src="data:{escape_attr(content_type)};base64,{encoded}" '
        f'alt="{escape_attr(alt)}">'
    )


def _render_table(table: Any) -> str:
    parts = ['<table class="pptx-table">']
    for r, row in enumerate(table.rows):
        parts.append("<tr>")
        cell_tag = "th" if r == 0 else "td"
        for cell in row.cells:
            parts.append(
                f"<{cell_tag}>{escape_text(cell.text).replace(chr(10), '<br>')}</{cell_tag}>"
            )
        parts.append("</tr>")
    parts.append("</table>")
    return "".join(parts)


def _shape_style(shape: Any, slide_width: int, slide_height: int) -> str:
    left = _pct(_emu(shape.left), slide_width)
    top = _pct(_emu(shape.top), slide_height)
    width = _pct(_emu(shape.width), slide_width)
    height = _pct(_emu(shape.height), slide_height)
    return f"left:{left};top:{top};width:{width};height:{height};"


def _emu(value: Any) -> int:
    return int(value) if value is not None else 0


def _pct(value: int, total: int) -> str:
    if total <= 0:
        return "0%"
    return f"{(value / total) * 100:.4f}%"


def _shape_title(shape: Any) -> str | None:
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            return text.splitlines()[0]
    return None
